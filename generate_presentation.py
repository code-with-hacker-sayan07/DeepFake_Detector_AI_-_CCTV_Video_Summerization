import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------------------
    # Styling Constants
    # -------------------------------------------------------------------------
    # Color Palette: Cyberpunk Dark Mode
    BG_COLOR = RGBColor(11, 15, 25)          # 0B0F19 (Deep Charcoal)
    CARD_BG = RGBColor(18, 24, 38)            # 121826 (Sleek Dark Navy)
    CARD_BORDER = RGBColor(30, 41, 59)        # 1E293B (Slate Border)
    
    CYAN = RGBColor(6, 182, 212)             # 06B6D4 (Cyber Cyan)
    GREEN = RGBColor(16, 185, 129)           # 10B981 (Neon Green)
    ORANGE = RGBColor(249, 115, 22)          # F97316 (Alert Orange)
    RED = RGBColor(239, 68, 68)              # EF4444 (Threat Red)
    
    WHITE = RGBColor(255, 255, 255)
    TEXT_LIGHT = RGBColor(226, 232, 240)     # E2E8F0 (Off-white body text)
    TEXT_GRAY = RGBColor(148, 163, 184)      # 94A3B8 (Slate Gray)

    FONT_FAMILY = "Segoe UI"

    # Helper function: set slide background to dark
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Helper function: add a professional header to content slides
    def add_header(slide, title, category="SENTRY-AI // SYSTEM ANALYSIS"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_right = cat_tf.margin_top = cat_tf.margin_bottom = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category.upper()
        cat_p.font.bold = True
        cat_p.font.size = Pt(10)
        cat_p.font.color.rgb = CYAN
        cat_p.font.name = FONT_FAMILY
        
        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_right = title_tf.margin_top = title_tf.margin_bottom = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.bold = True
        title_p.font.size = Pt(28)
        title_p.font.color.rgb = WHITE
        title_p.font.name = FONT_FAMILY
        
        # Cyan accent line under header
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = CYAN
        line.line.fill.background() # No border

    # Helper function: add card box layout
    def add_card(slide, left, top, width, height, card_title, bullets, accent_color=CYAN):
        # Background Card Shape
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)
        
        # Text Frame
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.25)
        tf.margin_bottom = Inches(0.25)
        
        # Title of the Card
        title_p = tf.paragraphs[0]
        title_p.text = "■  " + card_title.upper()
        title_p.font.bold = True
        title_p.font.size = Pt(15)
        title_p.font.color.rgb = accent_color
        title_p.font.name = FONT_FAMILY
        title_p.space_after = Pt(14)
        
        # Bullet Points
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(11.5)
            p.font.color.rgb = TEXT_LIGHT
            p.font.name = FONT_FAMILY
            p.space_after = Pt(8)
            p.level = 0
            
            # Simple manual formatting for clean hierarchy:
            # Let's replace simple dash bullets with a custom indented text
            p.line_spacing = 1.15

    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    # Subtle accent grid/line decorations to look cybernetic
    top_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(0.02))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = CARD_BORDER
    top_line.line.fill.background()
    
    bot_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.7), Inches(11.733), Inches(0.02))
    bot_line.fill.solid()
    bot_line.fill.fore_color.rgb = CARD_BORDER
    bot_line.line.fill.background()

    # Small top tags
    tag_left = slide1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(4.0), Inches(0.3))
    tag_left.text_frame.margin_left = tag_left.text_frame.margin_top = 0
    tl_p = tag_left.text_frame.paragraphs[0]
    tl_p.text = "VISION FORENSICS GROUP // TECHNICAL OVERVIEW"
    tl_p.font.size = Pt(9)
    tl_p.font.bold = True
    tl_p.font.color.rgb = CYAN
    tl_p.font.name = FONT_FAMILY

    tag_right = slide1.shapes.add_textbox(Inches(8.5), Inches(0.4), Inches(4.0), Inches(0.3))
    tag_right.text_frame.margin_left = tag_right.text_frame.margin_top = 0
    tr_p = tag_right.text_frame.paragraphs[0]
    tr_p.text = "CLASSIFICATION: INTERNAL USE ONLY"
    tr_p.font.size = Pt(9)
    tr_p.font.bold = True
    tr_p.font.color.rgb = RED
    tr_p.font.name = FONT_FAMILY
    tr_p.alignment = PP_ALIGN.RIGHT

    # Main Title
    main_title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(1.5))
    main_tf = main_title_box.text_frame
    main_tf.word_wrap = True
    main_p = main_tf.paragraphs[0]
    main_p.text = "S E N T R Y - A I"
    main_p.alignment = PP_ALIGN.CENTER
    main_p.font.bold = True
    main_p.font.size = Pt(60)
    main_p.font.color.rgb = CYAN
    main_p.font.name = FONT_FAMILY
    
    # Subtitle Line
    sub_title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(11.333), Inches(1.0))
    sub_tf = sub_title_box.text_frame
    sub_tf.word_wrap = True
    sub_p = sub_tf.paragraphs[0]
    sub_p.text = "Next-Generation DeepForensics & Asynchronous CCTV Intelligence Platform"
    sub_p.alignment = PP_ALIGN.CENTER
    sub_p.font.bold = True
    sub_p.font.size = Pt(20)
    sub_p.font.color.rgb = WHITE
    sub_p.font.name = FONT_FAMILY

    # Cyber Cyan Divider Bar
    div_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.666), Inches(4.5), Inches(4.0), Inches(0.05))
    div_bar.fill.solid()
    div_bar.fill.fore_color.rgb = CYAN
    div_bar.line.fill.background()

    # Meta Info
    meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.333), Inches(1.0))
    meta_tf = meta_box.text_frame
    meta_tf.word_wrap = True
    meta_p = meta_tf.paragraphs[0]
    meta_p.text = "Architectural Design, Algorithms, and Multi-Dimensional Engine Specifications"
    meta_p.alignment = PP_ALIGN.CENTER
    meta_p.font.size = Pt(13)
    meta_p.font.color.rgb = TEXT_GRAY
    meta_p.font.name = FONT_FAMILY

    # =========================================================================
    # SLIDE 2: The Threat Landscape
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "The Threat Landscape: Generative AI Manipulation", "SENTRY-AI // SECURING VISION INTEGRITY")
    
    add_card(
        slide=slide2,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="The Core Security Challenges",
        bullets=[
            "Rapid evolution of Generative Adversarial Networks (GANs) and Diffusion models allows threat actors to create highly realistic deepfakes.",
            "Visual evidence integrity is compromised, undermining legal evidence, identity validation, and corporate communications.",
            "Manual inspection of hours of CCTV and media feeds is subjective, exhausting, and fails to identify sub-pixel mathematical manipulations.",
            "Modern threat detection requires a unified, high-throughput solution that combines media forensics and motion intelligence."
        ],
        accent_color=RED
    )
    
    add_card(
        slide=slide2,
        left=Inches(6.9),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="The SENTRY-AI Response",
        bullets=[
            "Dual-Engine Defense: A unified framework implementing high-dimensional Deepfake Forensics and automated CCTV Summarization.",
            "Multi-Modal Verification: Checks for structural anomalies, frequency domain traces, and temporal optical flow inconsistencies.",
            "Automated Footage Compression: Clusters and isolates critical activity timestamps, reducing inspection times by over 80%.",
            "Operator-First Intelligence: Provides immediate executive briefs, anomaly grids, and threat scoring in a premium React command console."
        ],
        accent_color=GREEN
    )

    # =========================================================================
    # SLIDE 3: System Architecture
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "System Architecture & Core Technology Stack", "SENTRY-AI // DEPLOYMENT SCHEMA")
    
    add_card(
        slide=slide3,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="High-Performance Tech Stack",
        bullets=[
            "FastAPI Backend Portal: Asynchronous file processing running Python 3.14 with structured request validation schemas.",
            "Vision & Neural Foundations: Leverages PyTorch deep learning, Torchvision neural detectors, and highly optimized OpenCV pipelines.",
            "Modern Forensics Console: React 18 frontend built with Vite, styled with Tailwind CSS, and equipped with responsive Recharts graphs.",
            "Zero-Leak Memory Engine: Implements background session-based file storage management with automatic age-based cleanups."
        ],
        accent_color=CYAN
    )
    
    add_card(
        slide=slide3,
        left=Inches(6.9),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Modular Processing Pipeline",
        bullets=[
            "1. Media Ingestion: Handles video/image payloads up to 100MB with extensions validation.",
            "2. Uniform Sampling: Extracts exact frame intervals (e.g. 30 frames) for rapid forensic inspection without losing sequence context.",
            "3. Dual Engine Routing: Routes payloads dynamically to the Forensic Engine (deepfakes) or CCTV Summarizer Engine (footage intelligence).",
            "4. Dynamic Report Compilation: Combines metrics into unified JSON reports for client-side rendering and local analytics storage."
        ],
        accent_color=GREEN
    )

    # =========================================================================
    # SLIDE 4: Deepfake Engine - Spatial Domain
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Deepfake Forensic Engine: Spatial Domain Analysis", "SENTRY-AI // SPATIAL DETECTORS")
    
    add_card(
        slide=slide4,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Boundary Blending Degradation",
        bullets=[
            "Manipulated faces must be blended/warped onto target bodies. This leaves telltale anomalies at the face splice margins.",
            "The engine isolates detected face crops (via MTCNN or Haar Cascades) and segments them into border and inner zones.",
            "It computes the Laplacian variance ratio of boundary vs inner area: natural faces have progressive ratios (~0.7-1.1).",
            "Synthetically warped borders show heavily smoothed borders (low ratio) or excessive sharpening artifacts (high ratio)."
        ],
        accent_color=CYAN
    )
    
    add_card(
        slide=slide4,
        left=Inches(6.9),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Lighting Asymmetry & Texture Noise",
        bullets=[
            "Lighting Vector Asymmetry: Compares horizontal luminance gradients (Sobel X filters) between left and right facial halves.",
            "Splice detection catches inconsistencies where artificial lighting angles on the face mismatch the scene background.",
            "Texture Noise Coherence: Runs standard deviation checks across localized 8x8 pixels to track generative noise profiles.",
            "Symmetry Verification: Measures eye structure, sizes, and relative distances to flag structural facial asymmetries."
        ],
        accent_color=CYAN
    )

    # =========================================================================
    # SLIDE 5: Deepfake Engine - Frequency & Temporal
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Deepfake Forensic Engine: Spectral & Temporal Domains", "SENTRY-AI // FREQUENCY & TEMPORAL")
    
    add_card(
        slide=slide5,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="2D FFT Spectral Analysis",
        bullets=[
            "Generative algorithms (GANs, Diffusion) utilize upsampling (deconvolution) layers, which leave regular artifacts in the frequency domain.",
            "The engine applies a 2D Fast Fourier Transform (FFT) on face crops to compute their radial magnitude spectrum profiles.",
            "It separates low frequencies (overall face structure) from high frequencies (fine textures) to measure relative energy.",
            "Spikes or flat high-frequency profiles expose generative signatures. It outputs a 10x10 local spectral anomaly grid."
        ],
        accent_color=ORANGE
    )
    
    add_card(
        slide=slide5,
        left=Inches(6.9),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Spatiotemporal Optical Flow",
        bullets=[
            "Deepfakes struggle with temporal continuity: synthetic faces often 'flutter' or warp slightly from frame to frame.",
            "The temporal detector tracks consecutive sampled frames and crops the primary face region.",
            "It calculates Farneback Dense Optical Flow between frames to monitor flow magnitude, velocity standard deviation, and peak shifts.",
            "Natural facial movement is smooth and continuous. High-frequency velocity fluctuations indicate temporal anomalies."
        ],
        accent_color=ORANGE
    )

    # =========================================================================
    # SLIDE 6: CCTV Summarizer
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "CCTV Summarizer: Motion Profiling & Compression", "SENTRY-AI // footage compression")
    
    add_card(
        slide=slide6,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Active Motion Profile Extraction",
        bullets=[
            "To process CCTV footage efficiently, the engine monitors raw motion changes across sampled frames.",
            "It converts frames to grayscale, applies a Gaussian Blur (15x15) to reduce noise, and calculates absolute pixel differences.",
            "Applying binary thresholds and dilations isolates cohesive moving bodies while removing background compression noise.",
            "It tracks the active pixel ratio to generate a visual motion profile timeline for immediate charting."
        ],
        accent_color=CYAN
    )
    
    add_card(
        slide=slide6,
        left=Inches(6.9),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Timeline Activity Event Clustering",
        bullets=[
            "Groups contiguous frames exceeding a 1.2% motion threshold into discrete security 'Activity Events'.",
            "Coalesces active events separated by less than 2.0 seconds to keep activities (e.g. crossing a yard) in a single clip.",
            "Isolates the peak motion timestamp in each event to select a high-fidelity keyframe.",
            "Allows operators to instantly skip blank periods, review summaries, and cut security footage inspection time."
        ],
        accent_color=GREEN
    )

    # =========================================================================
    # SLIDE 7: CCTV Neural Vision
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "CCTV Neural Vision: Object Tracking & Threats", "SENTRY-AI // NEURAL DETECTION")
    
    add_card(
        slide=slide7,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="SSDLite MobileNetV3 Tracking",
        bullets=[
            "Runs a lightweight neural object detector (SSDLite MobileNetV3) tracking 15 COCO classes (Person, Vehicle, Bags, etc.).",
            "Maintains high speed (~12MB model size) suitable for local dev servers and real-time security systems.",
            "Heuristic Fallback: Uses aspect-ratio contours (tall for persons, wide for vehicles) if hardware constraints disable PyTorch.",
            "Calculates dynamic threat scores combining category weights, motion velocity, and late-night parameters."
        ],
        accent_color=ORANGE
    )
    
    add_card(
        slide=slide7,
        left=Inches(6.9),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Heuristic Security Violations",
        bullets=[
            "Night Intrusion: Flags a person detected during restricted night-time hours (9 PM - 6 AM) with immediate alerts.",
            "Unattended Baggage: Detects bags (backpack, suitcase) that remain stationary and unattended for more than 5 seconds.",
            "Reckless Driving: Triggers if a vehicle velocity (pixel-shift delta) exceeds extreme safety thresholds.",
            "Erratic Pedestrian Activity: Monitors rapid motion standard deviations to identify physical altercations or panic running."
        ],
        accent_color=RED
    )

    # =========================================================================
    # SLIDE 8: Cyber Dashboard
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Interactive Cyber Dashboard: User Experience", "SENTRY-AI // COMMAND CONSOLE")
    
    add_card(
        slide=slide8,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Visual Forensic Console",
        bullets=[
            "Cyberpunk Aesthetic: Modern dashboard utilizing sleek neon-accented dark modes, clear layout cards, and custom CSS styling.",
            "Interactive Timelines: Operators can scroll through isolated security events, showing timestamps, duration, and labels.",
            "Keyframe Visualizer: Displays high-resolution full frames alongside cropped anomaly heatmaps for immediate verification.",
            "Active Motion Charts: Renders dynamic line graphs of motion profiles to visually flag peak velocity event timestamps."
        ],
        accent_color=CYAN
    )
    
    add_card(
        slide=slide8,
        left=Inches(6.9),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Operator Workflow Tools",
        bullets=[
            "Executive Intelligence Brief: Automatically generates a natural language overview summarizing events and metrics.",
            "Verdicts and Alert Banners: Uses bright status indicators (Secure - Green, Warning - Orange, Critical - Red) for alerts.",
            "Dual Modes Selector: Allows quick switching between Deepfake Forensic analysis and CCTV summarization modules.",
            "System Health Portal: Real-time API monitoring showing server latency, model status, and CUDA hardware availability."
        ],
        accent_color=GREEN
    )

    # =========================================================================
    # SLIDE 9: Project Summary
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Project Summary & Future Roadmap", "SENTRY-AI // MILESTONES & ROADMAP")
    
    add_card(
        slide=slide9,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Current System Milestones",
        bullets=[
            "Unified FastAPI API backend processing asynchronous video and image analysis payloads.",
            "Hybrid Forensic Engine (MTCNN, 2D FFT, Laplacian Border variance, Farneback Dense Flow, PyTorch classifier).",
            "Responsive React Frontend Console featuring real-time visual heatmaps, timeline events, and line graphs.",
            "PyTorch classifier training script (train.py) successfully mapped for local dataset compilation and weight export."
        ],
        accent_color=GREEN
    )
    
    add_card(
        slide=slide9,
        left=Inches(6.9),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        card_title="Future Enhancements",
        bullets=[
            "Custom YOLOv8 Integration: Replace SSDLite with customized YOLO models for precise detection in low-light conditions.",
            "Live RTSP Video Streaming: Integrate WebSockets to process active surveillance streams with sub-second latency.",
            "Decentralized Edge Deployment: Containerize applications using Docker for deployment on local edge hardware.",
            "Active Learning Pipeline: Auto-flag low-confidence verdicts for operator confirmation to continually train and refine weights."
        ],
        accent_color=CYAN
    )

    # -------------------------------------------------------------------------
    # Save the presentation
    # -------------------------------------------------------------------------
    output_filename = "SENTRY-AI_Platform_Presentation.pptx"
    prs.save(output_filename)
    print(f"[SUCCESS] PowerPoint presentation saved successfully to: {output_filename}")
    return output_filename

if __name__ == "__main__":
    create_presentation()
